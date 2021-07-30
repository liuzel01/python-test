from pptx import Presentation
prs = Presentation("测试.pptx")
for index, slide in enumerate(prs.slides):
    print(f"第 {index+1} 页")
    for shape in slide.shapes:
        if shape.has_text_frame:
            text_frame = shape.text_frame
            # print(text_frame.text)
            # 如果分段读就用下面的代码
            for paragraph in text_frame.paragraphs:
                print(paragraph.text)
