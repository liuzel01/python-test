
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches


prs = Presentation()

# 添加版式
SLD_LAYOUT_TITLE_AND_CONTENT = 1 ##标题和内容版式的序号
slide_layout = prs.slide_layouts[SLD_LAYOUT_TITLE_AND_CONTENT]
slide = prs.slides.add_slide(slide_layout)



# 添加自动形状

shapes = slide.shapes
left = top = width = height = Inches(1.0)


shape = shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
)


left = top = width = height = Inches(3.0)

shape = shapes.add_shape(
    MSO_SHAPE.LINE_CALLOUT_2_ACCENT_BAR, left, top, width, height
)

# get the callout line coming out of the right place
adjs = shape.adjustments
adjs[0] = 0.5   # vert pos of junction in margin line, 0 is top
adjs[1] = 0.0   # horz pos of margin ln wrt shape width, 0 is left side
adjs[2] = 0.5   # vert pos of elbow wrt margin line, 0 is top
adjs[3] = -0.1  # horz pos of elbow wrt shape width, 0 is margin line
adjs[4] = 3.0   # vert pos of line end wrt shape height, 0 is top
a5 = adjs[3] - (adjs[4] - adjs[0]) * height/width
adjs[5] = a5    # horz pos of elbow wrt shape width, 0 is margin line

# rotate 45 degrees counter-clockwise
shape.rotation = -45.0


prs.save('新建幻灯片.pptx')

